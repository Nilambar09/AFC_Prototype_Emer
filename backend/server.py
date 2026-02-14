from fastapi import FastAPI, APIRouter, HTTPException, Depends, UploadFile, File, Form, status
from fastapi.security import HTTPBearer, HTTPAuthorizationCredentials
from dotenv import load_dotenv
from starlette.middleware.cors import CORSMiddleware
from motor.motor_asyncio import AsyncIOMotorClient
import os
import logging
from pathlib import Path
from pydantic import BaseModel, Field, EmailStr, ConfigDict
from typing import List, Optional
import uuid
from datetime import datetime, timezone, timedelta
import jwt
import bcrypt
import aiofiles
import base64
import io
from PyPDF2 import PdfReader
from pptx import Presentation
from PIL import Image

ROOT_DIR = Path(__file__).parent
load_dotenv(ROOT_DIR / '.env')

# MongoDB connection
mongo_url = os.environ['MONGO_URL']
client = AsyncIOMotorClient(mongo_url)
db = client[os.environ['DB_NAME']]

# JWT Configuration
JWT_SECRET = os.environ.get('JWT_SECRET', 'ventur_secret_key')
JWT_ALGORITHM = "HS256"
JWT_EXPIRATION_HOURS = 24

# Upload directory
UPLOAD_DIR = Path(os.environ.get('UPLOAD_DIR', '/app/backend/uploads'))
UPLOAD_DIR.mkdir(parents=True, exist_ok=True)

# Create the main app without a prefix
app = FastAPI(title="Ventur API", description="AI Pitch Deck Analyzer & Data Room")

# Create a router with the /api prefix
api_router = APIRouter(prefix="/api")

# Security
security = HTTPBearer()

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# ==================== Models ====================

class UserRegister(BaseModel):
    email: EmailStr
    password: str
    name: str
    company_name: Optional[str] = None

class UserLogin(BaseModel):
    email: EmailStr
    password: str

class UserResponse(BaseModel):
    model_config = ConfigDict(extra="ignore")
    id: str
    email: str
    name: str
    company_name: Optional[str] = None
    created_at: str

class TokenResponse(BaseModel):
    access_token: str
    token_type: str = "bearer"
    user: UserResponse

class PitchDeckResponse(BaseModel):
    model_config = ConfigDict(extra="ignore")
    id: str
    user_id: str
    filename: str
    file_type: str
    file_size: int
    analysis: Optional[dict] = None
    status: str
    created_at: str
    updated_at: str

class DataRoomCategory:
    SUMMARY = "summary"
    FINANCIALS = "financials"
    LEGAL = "legal"
    PREVIOUS_FUNDING = "previous_funding"
    INTELLECTUAL_PROPERTY = "intellectual_property"
    STAFF = "staff"
    METRICS = "metrics"
    OTHER = "other"

DATA_ROOM_CATEGORIES = [
    {"value": "pitch_deck", "label": "Pitch Deck & Summary", "description": "Pitch Deck, Executive Summary, One Pager"},
    {"value": "financial_statements", "label": "Financial Statements", "description": "P&L Statement, Balance Sheet, Cash Flow"},
    {"value": "financial_model", "label": "Financial Model & Projections", "description": "Revenue Model, Budget, Runway Analysis"},
    {"value": "cap_table", "label": "Cap Table & Equity", "description": "Capitalization Table, Stock Agreements, Vesting"},
    {"value": "legal_corporate", "label": "Legal & Corporate Docs", "description": "Articles of Incorporation, Bylaws, Board Consents"},
    {"value": "funding_history", "label": "Funding History", "description": "Previous Round Docs, Investor Rights, SAFE/Convertibles"},
    {"value": "intellectual_property", "label": "Intellectual Property", "description": "Patents, Trademarks, Trade Secrets, Brand Assets"},
    {"value": "team_hr", "label": "Team & HR", "description": "Org Chart, Employee List, Key Contracts, Advisors"},
    {"value": "metrics_kpis", "label": "Metrics & KPIs", "description": "Sales Pipeline, MRR/ARR, User Growth, Churn"},
    {"value": "product_tech", "label": "Product & Technology", "description": "Product Roadmap, Architecture, API Docs, Security"}
]

class DataRoomDocumentResponse(BaseModel):
    model_config = ConfigDict(extra="ignore")
    id: str
    user_id: str
    filename: str
    file_type: str
    file_size: int
    category: str
    subcategory: Optional[str] = None
    analysis: Optional[dict] = None
    status: str
    created_at: str
    updated_at: str

class AnalysisRequest(BaseModel):
    document_id: str

class DashboardStats(BaseModel):
    total_pitch_decks: int
    total_documents: int
    analyzed_documents: int
    pending_analysis: int
    documents_by_category: dict

# ==================== Helpers ====================

def hash_password(password: str) -> str:
    return bcrypt.hashpw(password.encode('utf-8'), bcrypt.gensalt()).decode('utf-8')

def verify_password(password: str, hashed: str) -> bool:
    return bcrypt.checkpw(password.encode('utf-8'), hashed.encode('utf-8'))

def create_token(user_id: str, email: str) -> str:
    payload = {
        "user_id": user_id,
        "email": email,
        "exp": datetime.now(timezone.utc) + timedelta(hours=JWT_EXPIRATION_HOURS)
    }
    return jwt.encode(payload, JWT_SECRET, algorithm=JWT_ALGORITHM)

async def get_current_user(credentials: HTTPAuthorizationCredentials = Depends(security)):
    try:
        token = credentials.credentials
        payload = jwt.decode(token, JWT_SECRET, algorithms=[JWT_ALGORITHM])
        user = await db.users.find_one({"id": payload["user_id"]}, {"_id": 0})
        if not user:
            raise HTTPException(status_code=401, detail="User not found")
        return user
    except jwt.ExpiredSignatureError:
        raise HTTPException(status_code=401, detail="Token expired")
    except jwt.InvalidTokenError:
        raise HTTPException(status_code=401, detail="Invalid token")

def get_file_extension(filename: str) -> str:
    return filename.rsplit('.', 1)[-1].lower() if '.' in filename else ''

async def extract_text_from_pdf(file_path: Path) -> str:
    """Extract text from PDF file"""
    try:
        reader = PdfReader(str(file_path))
        text = ""
        for page in reader.pages:
            text += page.extract_text() or ""
        return text
    except Exception as e:
        logger.error(f"Error extracting text from PDF: {e}")
        return ""

async def extract_text_from_pptx(file_path: Path) -> str:
    """Extract text from PowerPoint file"""
    try:
        prs = Presentation(str(file_path))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        logger.error(f"Error extracting text from PPTX: {e}")
        return ""

async def extract_text_from_image(file_path: Path) -> str:
    """For images, return a placeholder as we'll send the image to AI"""
    return "[Image file - will be analyzed by AI]"

async def analyze_with_ai(content: str, file_path: Path, file_type: str, analysis_type: str = "pitch_deck") -> dict:
    """Analyze content using Gemini (supports file attachments)"""
    try:
        from emergentintegrations.llm.chat import LlmChat, UserMessage, FileContentWithMimeType
        
        api_key = os.environ.get('EMERGENT_LLM_KEY')
        if not api_key:
            return {"error": "AI API key not configured"}
        
        session_id = str(uuid.uuid4())
        
        if analysis_type == "pitch_deck":
            system_message = """You are an expert startup pitch deck consultant with experience reviewing thousands of pitch decks for Y Combinator, Sequoia, and other top VCs. 

Analyze the pitch deck and provide detailed, actionable feedback in the following JSON structure:
{
    "overall_score": (1-10),
    "executive_summary": "Brief summary of the pitch deck's strengths and weaknesses",
    "sections_analysis": [
        {
            "section": "Problem",
            "score": (1-10),
            "feedback": "Specific feedback",
            "improvements": ["List of specific improvements"],
            "example_rewrite": "Example of better copy if applicable"
        }
    ],
    "visual_recommendations": {
        "overall_design": "Feedback on design",
        "charts_needed": ["List of charts/visualizations that would strengthen the deck"],
        "images_to_add": ["Suggestions for images to add"],
        "images_to_remove": ["Any images that should be removed or replaced"]
    },
    "content_improvements": [
        {
            "original_text": "Original sentence or paragraph",
            "suggested_text": "Improved version",
            "reason": "Why this change improves the pitch"
        }
    ],
    "missing_elements": ["Key elements missing from the pitch deck"],
    "investor_perspective": "How an investor would likely perceive this pitch deck",
    "next_steps": ["Prioritized list of what to fix first"]
}"""
        else:
            system_message = """You are an expert startup due diligence consultant specializing in data room organization and investor readiness.

Analyze this document and provide detailed feedback in the following JSON structure:
{
    "document_type": "Identified type of document",
    "completeness_score": (1-10),
    "summary": "Brief summary of the document",
    "key_findings": ["Important information found in the document"],
    "missing_information": ["Information that should be included but is missing"],
    "red_flags": ["Any concerning elements investors might notice"],
    "improvements": [
        {
            "area": "Area to improve",
            "current_state": "How it currently is",
            "recommendation": "What should be done",
            "priority": "high/medium/low"
        }
    ],
    "data_visualization_suggestions": [
        {
            "chart_type": "pie/bar/line/table",
            "data_to_visualize": "What data should be shown",
            "title": "Suggested chart title"
        }
    ],
    "investor_readiness": "How ready this document is for investor review",
    "next_steps": ["Prioritized actions to improve this document"]
}"""

        chat = LlmChat(
            api_key=api_key,
            session_id=session_id,
            system_message=system_message
        ).with_model("gemini", "gemini-2.5-flash")
        
        # Determine mime type
        ext = file_path.suffix.lower()
        mime_types = {
            '.pdf': 'application/pdf',
            '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            '.ppt': 'application/vnd.ms-powerpoint',
            '.png': 'image/png',
            '.jpg': 'image/jpeg',
            '.jpeg': 'image/jpeg'
        }
        mime_type = mime_types.get(ext, 'application/octet-stream')
        
        # Create message with file attachment for Gemini
        file_content = FileContentWithMimeType(
            file_path=str(file_path),
            mime_type=mime_type
        )
        
        prompt = f"""Please analyze this {analysis_type.replace('_', ' ')} document thoroughly. 
        
Additional context from text extraction:
{content[:5000] if content else 'No text could be extracted from this file.'}

Provide your analysis in the JSON format specified. Be specific, actionable, and helpful."""

        user_message = UserMessage(
            text=prompt,
            file_contents=[file_content]
        )
        
        response = await chat.send_message(user_message)
        
        # Try to parse as JSON
        import json
        try:
            # Extract JSON from response
            response_text = response.strip()
            if response_text.startswith('```json'):
                response_text = response_text[7:]
            if response_text.startswith('```'):
                response_text = response_text[3:]
            if response_text.endswith('```'):
                response_text = response_text[:-3]
            
            analysis = json.loads(response_text.strip())
            return {"success": True, "analysis": analysis}
        except json.JSONDecodeError:
            return {"success": True, "analysis": {"raw_feedback": response}}
            
    except Exception as e:
        logger.error(f"AI Analysis error: {e}")
        return {"success": False, "error": str(e)}

# ==================== Auth Routes ====================

@api_router.post("/auth/register", response_model=TokenResponse)
async def register(user_data: UserRegister):
    # Check if user exists
    existing = await db.users.find_one({"email": user_data.email})
    if existing:
        raise HTTPException(status_code=400, detail="Email already registered")
    
    # Create user
    user_id = str(uuid.uuid4())
    user_doc = {
        "id": user_id,
        "email": user_data.email,
        "password": hash_password(user_data.password),
        "name": user_data.name,
        "company_name": user_data.company_name,
        "created_at": datetime.now(timezone.utc).isoformat()
    }
    
    await db.users.insert_one(user_doc)
    
    token = create_token(user_id, user_data.email)
    
    return TokenResponse(
        access_token=token,
        user=UserResponse(
            id=user_id,
            email=user_data.email,
            name=user_data.name,
            company_name=user_data.company_name,
            created_at=user_doc["created_at"]
        )
    )

@api_router.post("/auth/login", response_model=TokenResponse)
async def login(credentials: UserLogin):
    user = await db.users.find_one({"email": credentials.email}, {"_id": 0})
    if not user or not verify_password(credentials.password, user["password"]):
        raise HTTPException(status_code=401, detail="Invalid email or password")
    
    token = create_token(user["id"], user["email"])
    
    return TokenResponse(
        access_token=token,
        user=UserResponse(
            id=user["id"],
            email=user["email"],
            name=user["name"],
            company_name=user.get("company_name"),
            created_at=user["created_at"]
        )
    )

@api_router.get("/auth/me", response_model=UserResponse)
async def get_me(current_user: dict = Depends(get_current_user)):
    return UserResponse(
        id=current_user["id"],
        email=current_user["email"],
        name=current_user["name"],
        company_name=current_user.get("company_name"),
        created_at=current_user["created_at"]
    )

# ==================== Pitch Deck Routes ====================

@api_router.post("/pitch-deck/upload", response_model=PitchDeckResponse)
async def upload_pitch_deck(
    file: UploadFile = File(...),
    current_user: dict = Depends(get_current_user)
):
    # Validate file type
    allowed_extensions = ['pdf', 'pptx', 'ppt', 'png', 'jpg', 'jpeg']
    ext = get_file_extension(file.filename)
    if ext not in allowed_extensions:
        raise HTTPException(status_code=400, detail=f"File type not allowed. Allowed: {', '.join(allowed_extensions)}")
    
    # Save file
    file_id = str(uuid.uuid4())
    file_path = UPLOAD_DIR / f"{file_id}.{ext}"
    
    content = await file.read()
    async with aiofiles.open(file_path, 'wb') as f:
        await f.write(content)
    
    # Create document record
    doc = {
        "id": file_id,
        "user_id": current_user["id"],
        "filename": file.filename,
        "file_type": ext,
        "file_size": len(content),
        "file_path": str(file_path),
        "analysis": None,
        "status": "uploaded",
        "created_at": datetime.now(timezone.utc).isoformat(),
        "updated_at": datetime.now(timezone.utc).isoformat()
    }
    
    await db.pitch_decks.insert_one(doc)
    
    return PitchDeckResponse(**{k: v for k, v in doc.items() if k != 'file_path'})

@api_router.post("/pitch-deck/{deck_id}/analyze", response_model=PitchDeckResponse)
async def analyze_pitch_deck(
    deck_id: str,
    current_user: dict = Depends(get_current_user)
):
    # Get pitch deck
    deck = await db.pitch_decks.find_one(
        {"id": deck_id, "user_id": current_user["id"]},
        {"_id": 0}
    )
    if not deck:
        raise HTTPException(status_code=404, detail="Pitch deck not found")
    
    # Update status
    await db.pitch_decks.update_one(
        {"id": deck_id},
        {"$set": {"status": "analyzing", "updated_at": datetime.now(timezone.utc).isoformat()}}
    )
    
    file_path = Path(deck["file_path"])
    
    # Extract text based on file type
    if deck["file_type"] == "pdf":
        content = await extract_text_from_pdf(file_path)
    elif deck["file_type"] in ["pptx", "ppt"]:
        content = await extract_text_from_pptx(file_path)
    else:
        content = await extract_text_from_image(file_path)
    
    # Analyze with AI
    analysis_result = await analyze_with_ai(content, file_path, deck["file_type"], "pitch_deck")
    
    # Update with analysis
    update_data = {
        "analysis": analysis_result,
        "status": "analyzed" if analysis_result.get("success") else "error",
        "updated_at": datetime.now(timezone.utc).isoformat()
    }
    
    await db.pitch_decks.update_one({"id": deck_id}, {"$set": update_data})
    
    updated_deck = await db.pitch_decks.find_one({"id": deck_id}, {"_id": 0, "file_path": 0})
    return PitchDeckResponse(**updated_deck)

@api_router.get("/pitch-decks", response_model=List[PitchDeckResponse])
async def get_pitch_decks(current_user: dict = Depends(get_current_user)):
    decks = await db.pitch_decks.find(
        {"user_id": current_user["id"]},
        {"_id": 0, "file_path": 0}
    ).sort("created_at", -1).to_list(100)
    return [PitchDeckResponse(**deck) for deck in decks]

@api_router.get("/pitch-deck/{deck_id}", response_model=PitchDeckResponse)
async def get_pitch_deck(deck_id: str, current_user: dict = Depends(get_current_user)):
    deck = await db.pitch_decks.find_one(
        {"id": deck_id, "user_id": current_user["id"]},
        {"_id": 0, "file_path": 0}
    )
    if not deck:
        raise HTTPException(status_code=404, detail="Pitch deck not found")
    return PitchDeckResponse(**deck)

@api_router.delete("/pitch-deck/{deck_id}")
async def delete_pitch_deck(deck_id: str, current_user: dict = Depends(get_current_user)):
    deck = await db.pitch_decks.find_one(
        {"id": deck_id, "user_id": current_user["id"]},
        {"_id": 0}
    )
    if not deck:
        raise HTTPException(status_code=404, detail="Pitch deck not found")
    
    # Delete file
    try:
        file_path = Path(deck["file_path"])
        if file_path.exists():
            file_path.unlink()
    except Exception as e:
        logger.error(f"Error deleting file: {e}")
    
    await db.pitch_decks.delete_one({"id": deck_id})
    return {"message": "Pitch deck deleted successfully"}

# ==================== Data Room Routes ====================

@api_router.get("/data-room/categories")
async def get_data_room_categories():
    return DATA_ROOM_CATEGORIES

@api_router.post("/data-room/upload", response_model=DataRoomDocumentResponse)
async def upload_data_room_document(
    file: UploadFile = File(...),
    category: str = Form(...),
    subcategory: Optional[str] = Form(None),
    current_user: dict = Depends(get_current_user)
):
    # Validate category
    valid_categories = [c["value"] for c in DATA_ROOM_CATEGORIES]
    if category not in valid_categories:
        raise HTTPException(status_code=400, detail="Invalid category")
    
    # Validate file type
    allowed_extensions = ['pdf', 'pptx', 'ppt', 'png', 'jpg', 'jpeg', 'xlsx', 'xls', 'doc', 'docx', 'csv']
    ext = get_file_extension(file.filename)
    if ext not in allowed_extensions:
        raise HTTPException(status_code=400, detail=f"File type not allowed. Allowed: {', '.join(allowed_extensions)}")
    
    # Save file
    file_id = str(uuid.uuid4())
    file_path = UPLOAD_DIR / f"dataroom_{file_id}.{ext}"
    
    content = await file.read()
    async with aiofiles.open(file_path, 'wb') as f:
        await f.write(content)
    
    # Create document record
    doc = {
        "id": file_id,
        "user_id": current_user["id"],
        "filename": file.filename,
        "file_type": ext,
        "file_size": len(content),
        "file_path": str(file_path),
        "category": category,
        "subcategory": subcategory,
        "analysis": None,
        "status": "uploaded",
        "created_at": datetime.now(timezone.utc).isoformat(),
        "updated_at": datetime.now(timezone.utc).isoformat()
    }
    
    await db.data_room_documents.insert_one(doc)
    
    return DataRoomDocumentResponse(**{k: v for k, v in doc.items() if k != 'file_path'})

@api_router.post("/data-room/{doc_id}/analyze", response_model=DataRoomDocumentResponse)
async def analyze_data_room_document(
    doc_id: str,
    current_user: dict = Depends(get_current_user)
):
    # Get document
    doc = await db.data_room_documents.find_one(
        {"id": doc_id, "user_id": current_user["id"]},
        {"_id": 0}
    )
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")
    
    # Update status
    await db.data_room_documents.update_one(
        {"id": doc_id},
        {"$set": {"status": "analyzing", "updated_at": datetime.now(timezone.utc).isoformat()}}
    )
    
    file_path = Path(doc["file_path"])
    
    # Extract text based on file type
    if doc["file_type"] == "pdf":
        content = await extract_text_from_pdf(file_path)
    elif doc["file_type"] in ["pptx", "ppt"]:
        content = await extract_text_from_pptx(file_path)
    else:
        content = await extract_text_from_image(file_path)
    
    # Analyze with AI
    analysis_result = await analyze_with_ai(content, file_path, doc["file_type"], f"data_room_{doc['category']}")
    
    # Update with analysis
    update_data = {
        "analysis": analysis_result,
        "status": "analyzed" if analysis_result.get("success") else "error",
        "updated_at": datetime.now(timezone.utc).isoformat()
    }
    
    await db.data_room_documents.update_one({"id": doc_id}, {"$set": update_data})
    
    updated_doc = await db.data_room_documents.find_one({"id": doc_id}, {"_id": 0, "file_path": 0})
    return DataRoomDocumentResponse(**updated_doc)

@api_router.get("/data-room", response_model=List[DataRoomDocumentResponse])
async def get_data_room_documents(
    category: Optional[str] = None,
    current_user: dict = Depends(get_current_user)
):
    query = {"user_id": current_user["id"]}
    if category:
        query["category"] = category
    
    docs = await db.data_room_documents.find(
        query,
        {"_id": 0, "file_path": 0}
    ).sort("created_at", -1).to_list(500)
    return [DataRoomDocumentResponse(**doc) for doc in docs]

@api_router.get("/data-room/{doc_id}", response_model=DataRoomDocumentResponse)
async def get_data_room_document(doc_id: str, current_user: dict = Depends(get_current_user)):
    doc = await db.data_room_documents.find_one(
        {"id": doc_id, "user_id": current_user["id"]},
        {"_id": 0, "file_path": 0}
    )
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")
    return DataRoomDocumentResponse(**doc)

@api_router.delete("/data-room/{doc_id}")
async def delete_data_room_document(doc_id: str, current_user: dict = Depends(get_current_user)):
    doc = await db.data_room_documents.find_one(
        {"id": doc_id, "user_id": current_user["id"]},
        {"_id": 0}
    )
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")
    
    # Delete file
    try:
        file_path = Path(doc["file_path"])
        if file_path.exists():
            file_path.unlink()
    except Exception as e:
        logger.error(f"Error deleting file: {e}")
    
    await db.data_room_documents.delete_one({"id": doc_id})
    return {"message": "Document deleted successfully"}

# ==================== Dashboard Routes ====================

@api_router.get("/dashboard/stats", response_model=DashboardStats)
async def get_dashboard_stats(current_user: dict = Depends(get_current_user)):
    user_id = current_user["id"]
    
    # Get pitch deck stats
    total_pitch_decks = await db.pitch_decks.count_documents({"user_id": user_id})
    
    # Get data room stats
    total_documents = await db.data_room_documents.count_documents({"user_id": user_id})
    analyzed_documents = await db.data_room_documents.count_documents({"user_id": user_id, "status": "analyzed"})
    
    # Get analyzed pitch decks
    analyzed_pitch_decks = await db.pitch_decks.count_documents({"user_id": user_id, "status": "analyzed"})
    
    # Get documents by category
    pipeline = [
        {"$match": {"user_id": user_id}},
        {"$group": {"_id": "$category", "count": {"$sum": 1}}}
    ]
    category_counts = await db.data_room_documents.aggregate(pipeline).to_list(100)
    documents_by_category = {item["_id"]: item["count"] for item in category_counts}
    
    pending_analysis = (total_pitch_decks - analyzed_pitch_decks) + (total_documents - analyzed_documents)
    
    return DashboardStats(
        total_pitch_decks=total_pitch_decks,
        total_documents=total_documents,
        analyzed_documents=analyzed_documents + analyzed_pitch_decks,
        pending_analysis=pending_analysis,
        documents_by_category=documents_by_category
    )

# ==================== History Routes ====================

@api_router.get("/history")
async def get_history(current_user: dict = Depends(get_current_user)):
    user_id = current_user["id"]
    
    # Get all pitch decks
    pitch_decks = await db.pitch_decks.find(
        {"user_id": user_id},
        {"_id": 0, "file_path": 0}
    ).sort("created_at", -1).to_list(100)
    
    # Get all data room documents
    data_room_docs = await db.data_room_documents.find(
        {"user_id": user_id},
        {"_id": 0, "file_path": 0}
    ).sort("created_at", -1).to_list(500)
    
    # Combine and format
    history = []
    
    for deck in pitch_decks:
        history.append({
            "id": deck["id"],
            "type": "pitch_deck",
            "filename": deck["filename"],
            "file_type": deck["file_type"],
            "status": deck["status"],
            "created_at": deck["created_at"],
            "has_analysis": deck.get("analysis") is not None
        })
    
    for doc in data_room_docs:
        history.append({
            "id": doc["id"],
            "type": "data_room",
            "filename": doc["filename"],
            "file_type": doc["file_type"],
            "category": doc["category"],
            "status": doc["status"],
            "created_at": doc["created_at"],
            "has_analysis": doc.get("analysis") is not None
        })
    
    # Sort by created_at
    history.sort(key=lambda x: x["created_at"], reverse=True)
    
    return history

@api_router.delete("/history/clear")
async def clear_history(current_user: dict = Depends(get_current_user)):
    user_id = current_user["id"]
    
    # Get all files to delete
    pitch_decks = await db.pitch_decks.find({"user_id": user_id}).to_list(1000)
    data_room_docs = await db.data_room_documents.find({"user_id": user_id}).to_list(1000)
    
    # Delete files
    for deck in pitch_decks:
        try:
            file_path = Path(deck.get("file_path", ""))
            if file_path.exists():
                file_path.unlink()
        except Exception as e:
            logger.error(f"Error deleting file: {e}")
    
    for doc in data_room_docs:
        try:
            file_path = Path(doc.get("file_path", ""))
            if file_path.exists():
                file_path.unlink()
        except Exception as e:
            logger.error(f"Error deleting file: {e}")
    
    # Delete from database
    await db.pitch_decks.delete_many({"user_id": user_id})
    await db.data_room_documents.delete_many({"user_id": user_id})
    
    return {"message": "History cleared successfully"}

# ==================== Health Check ====================

@api_router.get("/health")
async def health_check():
    return {"status": "healthy", "timestamp": datetime.now(timezone.utc).isoformat()}

# Include the router in the main app
app.include_router(api_router)

app.add_middleware(
    CORSMiddleware,
    allow_credentials=True,
    allow_origins=os.environ.get('CORS_ORIGINS', '*').split(','),
    allow_methods=["*"],
    allow_headers=["*"],
)

@app.on_event("shutdown")
async def shutdown_db_client():
    client.close()
